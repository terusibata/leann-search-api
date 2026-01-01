"""Document management API endpoints."""

import json
import uuid
from typing import Literal

from fastapi import APIRouter, File, Form, HTTPException, Query, UploadFile, status

from src.config import get_settings
from src.schemas.common import APIResponse, ErrorCode, PaginationInfo
from src.schemas.document import (
    AddDocumentsRequest,
    AddDocumentsResponse,
    BulkDeleteRequest,
    BulkDeleteResponse,
    DeleteDocumentResponse,
    DocumentDetailResponse,
    DocumentListResponse,
    FileUploadResponse,
    UpdateDocumentRequest,
    UpdateDocumentResponse,
    UpdateMetadataRequest,
    UpdateMetadataResponse,
)
from src.services.document_service import get_document_service
from src.services.index_service import get_index_service

router = APIRouter(prefix="/indexes/{index_name}/documents", tags=["documents"])


def _check_index_exists(index_name: str) -> None:
    """Check if index exists, raise 404 if not."""
    service = get_index_service()
    if not service.index_exists(index_name):
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.INDEX_NOT_FOUND,
                    "message": f"Index '{index_name}' not found",
                },
            },
        )


@router.post("", response_model=APIResponse[AddDocumentsResponse])
async def add_documents(index_name: str, request: AddDocumentsRequest):
    """Add documents to an index."""
    _check_index_exists(index_name)

    service = get_document_service()

    documents = [doc.model_dump() for doc in request.documents]
    options = request.options or {}

    results = service.add_documents(
        index_name,
        documents,
        chunk_size=options.chunk_size if hasattr(options, "chunk_size") else None,
        chunk_overlap=options.chunk_overlap if hasattr(options, "chunk_overlap") else None,
        update_if_exists=options.update_if_exists if hasattr(options, "update_if_exists") else False,
    )

    added = sum(1 for r in results if r.status == "added")
    updated = sum(1 for r in results if r.status == "updated")
    failed = sum(1 for r in results if r.status == "failed")

    return APIResponse.ok(
        AddDocumentsResponse(
            added=added,
            updated=updated,
            failed=failed,
            documents=results,
        )
    )


@router.post("/file", response_model=APIResponse[FileUploadResponse])
async def add_document_from_file(
    index_name: str,
    file: UploadFile = File(...),
    metadata: str | None = Form(None),
    document_id: str | None = Form(None),
):
    """Add a document from an uploaded file."""
    _check_index_exists(index_name)

    settings = get_settings()

    # Check file size
    content = await file.read()
    if len(content) > settings.max_upload_size_bytes:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.VALIDATION_ERROR,
                    "message": f"File size exceeds maximum allowed size of {settings.max_upload_size_mb}MB",
                },
            },
        )

    # Determine file type
    filename = file.filename or "unknown"
    file_ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    supported_types = {"txt", "md", "pdf", "docx", "pptx"}
    if file_ext not in supported_types:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.VALIDATION_ERROR,
                    "message": f"Unsupported file type: {file_ext}. Supported: {', '.join(supported_types)}",
                },
            },
        )

    # Parse metadata
    meta_dict = {}
    if metadata:
        try:
            meta_dict = json.loads(metadata)
        except json.JSONDecodeError:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail={
                    "success": False,
                    "data": None,
                    "error": {
                        "code": ErrorCode.VALIDATION_ERROR,
                        "message": "Invalid metadata JSON",
                    },
                },
            )

    # Extract text based on file type
    text_content = ""
    page_count = None

    if file_ext in ("txt", "md"):
        text_content = content.decode("utf-8", errors="ignore")
    elif file_ext == "pdf":
        try:
            import io

            import fitz  # PyMuPDF

            pdf_doc = fitz.open(stream=content, filetype="pdf")
            page_count = len(pdf_doc)
            text_content = "\n\n".join(page.get_text() for page in pdf_doc)
            pdf_doc.close()
        except ImportError:
            # Fallback: just try to decode as text
            text_content = content.decode("utf-8", errors="ignore")
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail={
                    "success": False,
                    "data": None,
                    "error": {
                        "code": ErrorCode.INTERNAL_ERROR,
                        "message": f"Failed to parse PDF: {str(e)}",
                    },
                },
            )
    elif file_ext == "docx":
        try:
            import io

            from docx import Document

            doc = Document(io.BytesIO(content))
            text_content = "\n\n".join(para.text for para in doc.paragraphs)
        except ImportError:
            text_content = content.decode("utf-8", errors="ignore")
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail={
                    "success": False,
                    "data": None,
                    "error": {
                        "code": ErrorCode.INTERNAL_ERROR,
                        "message": f"Failed to parse DOCX: {str(e)}",
                    },
                },
            )
    elif file_ext == "pptx":
        try:
            import io

            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            text_content = "\n\n".join(texts)
            page_count = len(prs.slides)
        except ImportError:
            text_content = content.decode("utf-8", errors="ignore")
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail={
                    "success": False,
                    "data": None,
                    "error": {
                        "code": ErrorCode.INTERNAL_ERROR,
                        "message": f"Failed to parse PPTX: {str(e)}",
                    },
                },
            )

    if not text_content.strip():
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.VALIDATION_ERROR,
                    "message": "No text content could be extracted from the file",
                },
            },
        )

    # Add filename to metadata
    meta_dict["filename"] = filename
    meta_dict["file_type"] = file_ext

    doc_id = document_id or str(uuid.uuid4())

    service = get_document_service()
    results = service.add_documents(
        index_name,
        [{"id": doc_id, "content": text_content, "metadata": meta_dict}],
        update_if_exists=True,
    )

    if results and results[0].status != "failed":
        return APIResponse.ok(
            FileUploadResponse(
                id=doc_id,
                filename=filename,
                file_type=file_ext,
                page_count=page_count,
                chunk_count=results[0].chunk_count,
                status=results[0].status,
            )
        )
    else:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.INTERNAL_ERROR,
                    "message": results[0].error if results else "Failed to add document",
                },
            },
        )


@router.get("", response_model=APIResponse[DocumentListResponse])
async def list_documents(
    index_name: str,
    page: int = Query(1, ge=1),
    per_page: int = Query(50, ge=1, le=100),
    sort_by: Literal["created_at", "updated_at", "id"] = "created_at",
    sort_order: Literal["asc", "desc"] = "desc",
    metadata_filter: str | None = Query(None),
):
    """List documents in an index."""
    _check_index_exists(index_name)

    # Parse metadata filter
    filter_dict = None
    if metadata_filter:
        try:
            filter_dict = json.loads(metadata_filter)
        except json.JSONDecodeError:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail={
                    "success": False,
                    "data": None,
                    "error": {
                        "code": ErrorCode.VALIDATION_ERROR,
                        "message": "Invalid metadata_filter JSON",
                    },
                },
            )

    service = get_document_service()
    documents, total = service.list_documents(
        index_name,
        page=page,
        per_page=per_page,
        sort_by=sort_by,
        sort_order=sort_order,
        metadata_filter=filter_dict,
    )

    total_pages = (total + per_page - 1) // per_page

    return APIResponse.ok(
        DocumentListResponse(
            documents=documents,
            pagination=PaginationInfo(
                page=page,
                per_page=per_page,
                total=total,
                total_pages=total_pages,
            ),
        )
    )


@router.get("/{document_id}", response_model=APIResponse[DocumentDetailResponse])
async def get_document(index_name: str, document_id: str):
    """Get document details."""
    _check_index_exists(index_name)

    service = get_document_service()
    document = service.get_document(index_name, document_id)

    if not document:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.DOCUMENT_NOT_FOUND,
                    "message": f"Document '{document_id}' not found",
                },
            },
        )

    return APIResponse.ok(document)


@router.put("/{document_id}", response_model=APIResponse[UpdateDocumentResponse])
async def update_document(
    index_name: str, document_id: str, request: UpdateDocumentRequest
):
    """Update a document."""
    _check_index_exists(index_name)

    service = get_document_service()

    if not service.document_exists(index_name, document_id):
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.DOCUMENT_NOT_FOUND,
                    "message": f"Document '{document_id}' not found",
                },
            },
        )

    merge_metadata = True
    if request.options and hasattr(request.options, "merge_metadata"):
        merge_metadata = request.options.merge_metadata

    result = service.update_document(
        index_name,
        document_id,
        content=request.content,
        metadata=request.metadata,
        merge_metadata=merge_metadata,
    )

    if result:
        return APIResponse.ok(
            UpdateDocumentResponse(
                id=result.id,
                chunk_count=result.chunk_count,
                status=result.status,
            )
        )
    else:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.INTERNAL_ERROR,
                    "message": "Failed to update document",
                },
            },
        )


@router.patch("/{document_id}/metadata", response_model=APIResponse[UpdateMetadataResponse])
async def update_document_metadata(
    index_name: str, document_id: str, request: UpdateMetadataRequest
):
    """Update only the metadata of a document (no re-indexing needed)."""
    _check_index_exists(index_name)

    service = get_document_service()

    if not service.document_exists(index_name, document_id):
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.DOCUMENT_NOT_FOUND,
                    "message": f"Document '{document_id}' not found",
                },
            },
        )

    new_metadata = service.update_metadata_only(
        index_name, document_id, request.metadata, merge=request.merge
    )

    if new_metadata is not None:
        return APIResponse.ok(
            UpdateMetadataResponse(id=document_id, metadata=new_metadata)
        )
    else:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.INTERNAL_ERROR,
                    "message": "Failed to update metadata",
                },
            },
        )


@router.delete("/{document_id}", response_model=APIResponse[DeleteDocumentResponse])
async def delete_document(index_name: str, document_id: str):
    """Delete a document."""
    _check_index_exists(index_name)

    service = get_document_service()

    if not service.document_exists(index_name, document_id):
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.DOCUMENT_NOT_FOUND,
                    "message": f"Document '{document_id}' not found",
                },
            },
        )

    service.delete_document(index_name, document_id)

    return APIResponse.ok(
        DeleteDocumentResponse(message=f"Document '{document_id}' deleted successfully")
    )


@router.post("/bulk-delete", response_model=APIResponse[BulkDeleteResponse])
async def bulk_delete_documents(index_name: str, request: BulkDeleteRequest):
    """Delete multiple documents."""
    _check_index_exists(index_name)

    if not request.document_ids and not request.metadata_filter:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail={
                "success": False,
                "data": None,
                "error": {
                    "code": ErrorCode.VALIDATION_ERROR,
                    "message": "Either document_ids or metadata_filter must be provided",
                },
            },
        )

    service = get_document_service()
    deleted = service.bulk_delete(
        index_name,
        document_ids=request.document_ids,
        metadata_filter=request.metadata_filter,
    )

    return APIResponse.ok(
        BulkDeleteResponse(
            deleted=deleted, message=f"Deleted {deleted} document(s) successfully"
        )
    )
